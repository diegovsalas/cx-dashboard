"""Genera presentación PPT profesional — Bot WhatsApp & Dashboard CX"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ═══════════════ COLORES ═══════════════
CHARCOAL = RGBColor(0x1E, 0x1E, 0x2E)
DARK = RGBColor(0x2D, 0x2D, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xF8, 0xF8, 0xFA)
BLUE = RGBColor(0x0D, 0x8B, 0xF2)
BLUE_SOFT = RGBColor(0xE3, 0xF0, 0xFD)
GREEN = RGBColor(0x25, 0xD3, 0x66)
GREEN_SOFT = RGBColor(0xE6, 0xF9, 0xEC)
RED = RGBColor(0xEF, 0x44, 0x44)
RED_SOFT = RGBColor(0xFD, 0xEC, 0xEC)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_SOFT = RGBColor(0xFE, 0xF5, 0xE0)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
GRAY = RGBColor(0x9C, 0x9C, 0xAC)
GRAY_DARK = RGBColor(0x6B, 0x6B, 0x7B)
BORDER = RGBColor(0xE2, 0xE2, 0xEA)
BG_CARD = RGBColor(0xF4, 0xF4, 0xF8)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ═══════════════ HELPERS ═══════════════
def rect(slide, l, t, w, h, fill, border=None, radius=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def txt(slide, l, t, w, h, text, sz=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.bold = bold; p.font.color.rgb = color; p.font.name = 'Calibri'; p.alignment = align
    return tb

def multi(slide, l, t, w, h, lines, sz=11, color=DARK, spacing=2):
    """lines: list of (text, bold, color_override)"""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, b, c) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(sz); p.font.bold = b
        p.font.color.rgb = c or color; p.font.name = 'Calibri'
        p.space_before = Pt(spacing)
    return tb

def pill(slide, l, t, text, fill, text_color):
    w = max(1.2, len(text) * 0.09 + 0.4)
    rect(slide, l, t, w, 0.3, fill, radius=True)
    txt(slide, l, t, w, 0.3, text, 9, True, text_color, PP_ALIGN.CENTER)
    return w

def line_h(slide, l, t, w, color=BORDER):
    rect(slide, l, t, w, 0.02, color)

# ═══════════════════════════════════════
# SLIDE 1: PORTADA
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, 13.33, 7.5, CHARCOAL)
rect(s, 0, 7.3, 13.33, 0.2, BLUE)

# Acento lateral
rect(s, 0, 0, 0.08, 7.5, BLUE)

txt(s, 1.5, 1.5, 4, 0.4, 'GRUPO AVANTEX', 13, True, GRAY, PP_ALIGN.LEFT)
line_h(s, 1.5, 1.95, 2, BLUE)

txt(s, 1.5, 2.3, 10, 1, 'Bot WhatsApp &\nDashboard CX', 44, True, WHITE)
txt(s, 1.5, 3.8, 7, 0.6, 'Nuevo sistema de atención al cliente para reportes\nde incidencias, quejas y cancelaciones de servicio', 16, False, GRAY)

# Badges
x = 1.5
for label, bg, tc in [('WhatsApp Bot', RGBColor(0x1A,0x3A,0x1E), GREEN), ('Dashboard CX', RGBColor(0x12,0x2A,0x45), BLUE), ('Chat Agentes', RGBColor(0x25,0x1A,0x40), PURPLE)]:
    w = pill(s, x, 4.7, label, bg, tc); x += w + 0.15

# WhatsApp number box
rect(s, 1.5, 5.5, 4.5, 0.65, RGBColor(0x15,0x2E,0x1A), GREEN, True)
txt(s, 1.5, 5.5, 4.5, 0.65, '+52 449 156 1238', 22, True, GREEN, PP_ALIGN.CENTER)

txt(s, 1.5, 6.7, 5, 0.3, 'Customer Experience  ·  Abril 2026', 10, False, GRAY)

# Decorative circles
for cx, cy, cr, cc in [(11.5, 2, 1.8, RGBColor(0x25,0x2A,0x35)), (10.8, 4.5, 1.2, RGBColor(0x22,0x27,0x32)), (12, 5.5, 0.8, RGBColor(0x28,0x2D,0x38))]:
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(cr), Inches(cr))
    sh.fill.solid(); sh.fill.fore_color.rgb = cc; sh.line.fill.background()

# ═══════════════════════════════════════
# SLIDE 2: FLUJO
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, 13.33, 7.5, OFFWHITE)
rect(s, 0, 0, 13.33, 0.06, BLUE)

txt(s, 0.8, 0.4, 1, 0.4, '01', 28, True, BLUE)
txt(s, 1.5, 0.45, 5, 0.4, 'Flujo del Bot', 24, True, CHARCOAL)
txt(s, 1.5, 0.85, 8, 0.3, 'Cómo funciona el reporte de incidencias y cancelaciones por WhatsApp', 11, False, GRAY_DARK)

# — Card: Incidencias —
rect(s, 0.5, 1.4, 5.9, 2.7, WHITE, BORDER, True)
rect(s, 0.5, 1.4, 5.9, 0.06, GREEN)
txt(s, 0.8, 1.55, 4, 0.3, '📋  Flujo de Incidencias', 13, True, CHARCOAL)

steps = ['Hola','Nombre','Empresa','Sucursal','Servicio','Tipo','Descripción','FOLIO ✅']
x = 0.8
for i, step in enumerate(steps):
    is_last = i == len(steps)-1
    bg = BLUE if is_last else BG_CARD
    tc = WHITE if is_last else CHARCOAL
    rect(s, x, 2.05, 0.58 if not is_last else 0.7, 0.28, bg, BORDER if not is_last else None, True)
    txt(s, x, 2.05, 0.58 if not is_last else 0.7, 0.28, step, 7, True, tc, PP_ALIGN.CENTER)
    x += (0.58 if not is_last else 0.7) + 0.04
    if not is_last and i < len(steps)-1:
        txt(s, x - 0.12, 2.03, 0.1, 0.3, '→', 9, False, GRAY)

multi(s, 0.8, 2.5, 5.3, 0.9, [
    ('Servicios:  Aromatex  ·  Pestex  ·  Weldex', True, BLUE),
    ('El cliente puede consultar su ticket, enviar evidencia', False, GRAY_DARK),
    ('o solicitar agente humano en cualquier momento.', False, GRAY_DARK),
], 10, GRAY_DARK, 1)

# — Card: Cancelaciones —
rect(s, 0.5, 4.3, 5.9, 2.7, WHITE, BORDER, True)
rect(s, 0.5, 4.3, 5.9, 0.06, RED)
txt(s, 0.8, 4.45, 4, 0.3, '🔴  Flujo de Cancelaciones', 13, True, CHARCOAL)

steps_c = ['"Cancelar"','Confirmar','Nombre','Empresa','Servicio','Motivo','FOLIO 🔴']
x = 0.8
for i, step in enumerate(steps_c):
    is_last = i == len(steps_c)-1
    bg = RED if is_last else BG_CARD
    tc = WHITE if is_last else CHARCOAL
    w = 0.65 if i==0 else (0.7 if is_last else 0.58)
    rect(s, x, 4.95, w, 0.28, bg, BORDER if not is_last else None, True)
    txt(s, x, 4.95, w, 0.28, step, 7, True, tc, PP_ALIGN.CENTER)
    x += w + 0.04
    if not is_last: txt(s, x - 0.12, 4.93, 0.1, 0.3, '→', 9, False, GRAY)

multi(s, 0.8, 5.45, 5.3, 1.2, [
    ('Detección automática de intención de cancelar.', False, GRAY_DARK),
    ('Se envían cláusulas de carta convenio.', False, GRAY_DARK),
    ('Prioridad Alta → asignado a nallelyquiroz.', True, RED),
], 10, GRAY_DARK, 1)

# — Card: Asignación —
rect(s, 6.8, 1.4, 6, 2.7, WHITE, BORDER, True)
rect(s, 6.8, 1.4, 6, 0.06, BLUE)
txt(s, 7.1, 1.55, 4, 0.3, '👥  Asignación Automática', 13, True, CHARCOAL)

agents = [
    ('Incidencias / Quejas', 'anakaren', BLUE),
    ('Cambio de aroma', 'anakaren', BLUE),
    ('Cliente pide agente', 'anakaren', BLUE),
    ('Cancelaciones', 'nallelyquiroz', RED),
    ('Escalado (kanban)', 'diegovelazquez', PURPLE),
]
for i, (tipo, agente, c) in enumerate(agents):
    y = 2.05 + i * 0.38
    txt(s, 7.1, y, 3.2, 0.3, f'→  {tipo}', 10, False, CHARCOAL)
    pill(s, 10.5, y + 0.03, agente, c, WHITE)

rect(s, 7.1, 3.65, 5.4, 0.3, BLUE_SOFT, radius=True)
txt(s, 7.3, 3.65, 5.2, 0.3, '🔔  Cada agente recibe WhatsApp con datos del cliente + link', 9, False, BLUE)

# — Card: Horario —
rect(s, 6.8, 4.3, 2.8, 2.7, WHITE, BORDER, True)
rect(s, 6.8, 4.3, 2.8, 0.06, AMBER)
txt(s, 7.1, 4.45, 2.5, 0.3, '⏰  Horario', 13, True, CHARCOAL)
multi(s, 7.1, 4.9, 2.5, 1.5, [
    ('Lunes a viernes', True, CHARCOAL),
    ('7:30 a.m. – 6:00 p.m.', False, GRAY_DARK),
    ('', False, None),
    ('Sábados', True, CHARCOAL),
    ('9:00 a.m. – 12:00 p.m.', False, GRAY_DARK),
    ('', False, None),
    ('Fuera de horario: ticket se', False, GRAY_DARK),
    ('registra → siguiente día hábil', False, GRAY_DARK),
], 9, GRAY_DARK, 1)

# — Card: CSAT —
rect(s, 10, 4.3, 2.8, 2.7, WHITE, BORDER, True)
rect(s, 10, 4.3, 2.8, 0.06, GREEN)
txt(s, 10.3, 4.45, 2.5, 0.3, '⭐  CSAT & NPS', 13, True, CHARCOAL)
multi(s, 10.3, 4.9, 2.3, 1.8, [
    ('Encuesta 1–5 al cerrar ticket', False, GRAY_DARK),
    ('', False, None),
    ('5  =  Promotor', True, GREEN),
    ('4  =  Pasivo', True, AMBER),
    ('1-3  =  Detractor', True, RED),
    ('', False, None),
    ('NPS = %Prom − %Det', True, CHARCOAL),
], 9, GRAY_DARK, 2)

# ═══════════════════════════════════════
# SLIDE 3: DASHBOARD
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, 13.33, 7.5, OFFWHITE)
rect(s, 0, 0, 13.33, 0.06, BLUE)

txt(s, 0.8, 0.4, 1, 0.4, '02', 28, True, BLUE)
txt(s, 1.5, 0.45, 8, 0.4, 'Dashboard CX & Chat de Agentes', 24, True, CHARCOAL)
txt(s, 1.5, 0.85, 8, 0.3, 'Plataformas para monitorear, gestionar y resolver tickets', 11, False, GRAY_DARK)

# — Dashboard CX —
rect(s, 0.5, 1.4, 6, 3, WHITE, BORDER, True)
rect(s, 0.5, 1.4, 6, 0.06, BLUE)
txt(s, 0.8, 1.55, 3, 0.3, '📊  Dashboard CX', 13, True, CHARCOAL)
txt(s, 4.2, 1.58, 2.5, 0.2, 'cx-dashboard-54yw.onrender.com', 7, False, GRAY)

secciones = [
    ('Inicio', 'Guía, WhatsApp, QR, descarga guía visual para clientes'),
    ('Dashboard', 'KPIs, gráficas tipo/servicio/zona/mes, NPS, CSAT'),
    ('Registros', 'Búsqueda global, filtros, modal de detalle editable'),
    ('Cancelaciones', 'Montos: Cancelado / Retenido / En negociación'),
    ('Tickets WA', 'Tickets del bot con zona editable, modal completo'),
]
for i, (sec, desc) in enumerate(secciones):
    y = 2.05 + i * 0.42
    rect(s, 0.8, y, 0.06, 0.25, BLUE)
    txt(s, 1.0, y, 1.3, 0.25, sec, 10, True, CHARCOAL)
    txt(s, 2.3, y, 4, 0.25, desc, 9, False, GRAY_DARK)

# — Chat —
rect(s, 0.5, 4.6, 6, 2.6, WHITE, BORDER, True)
rect(s, 0.5, 4.6, 6, 0.06, PURPLE)
txt(s, 0.8, 4.75, 3, 0.3, '💬  Chat de Agentes', 13, True, CHARCOAL)
txt(s, 4.2, 4.78, 2.5, 0.2, 'avantex-bot.onrender.com/dashboard', 7, False, GRAY)

chat_items = [
    ('Chat', 'Conversaciones en tiempo real con clientes'),
    ('Kanban', 'Nuevo → En proceso → Resuelto'),
    ('Stats', 'KPIs, NPS, CSAT en tiempo real'),
    ('Control', 'Agente toma conversación del bot'),
    ('CSAT', 'Enviar encuesta de satisfacción'),
]
for i, (label, desc) in enumerate(chat_items):
    y = 5.25 + i * 0.35
    rect(s, 0.8, y, 0.06, 0.22, PURPLE)
    txt(s, 1.0, y, 1, 0.22, label, 10, True, CHARCOAL)
    txt(s, 2.0, y, 4.3, 0.22, desc, 9, False, GRAY_DARK)

# — Roles —
rect(s, 6.8, 1.4, 6, 2.2, WHITE, BORDER, True)
rect(s, 6.8, 1.4, 6, 0.06, AMBER)
txt(s, 7.1, 1.55, 4, 0.3, '🔐  Roles y Permisos', 13, True, CHARCOAL)

roles = [
    ('Superadmin', 'diegovelazquez@', 'Todo + eliminar + editar nombres', RED),
    ('Admin', 'servicioalcliente@', 'Editar status, comentarios, montos, Chat', BLUE),
    ('Consulta', 'cobranza@, comercial@', 'Solo lectura en dashboards', AMBER),
]
for i, (rol, usr, perm, c) in enumerate(roles):
    y = 2.1 + i * 0.45
    pill(s, 7.1, y, rol, c, WHITE)
    txt(s, 8.6, y, 2, 0.25, usr, 8, False, GRAY)
    txt(s, 7.1, y + 0.23, 5.5, 0.2, perm, 8, False, GRAY_DARK)

# — Campos editables —
rect(s, 6.8, 3.8, 6, 3.4, WHITE, BORDER, True)
rect(s, 6.8, 3.8, 6, 0.06, GREEN)
txt(s, 7.1, 3.95, 5, 0.3, '✏️  Campos Editables', 13, True, CHARCOAL)

# 3 mini cards
for i, (title, items, c) in enumerate([
    ('Tickets', 'Status · Zona\nComentarios\nFecha solución', BLUE),
    ('Cancelaciones', 'Status cancelación\nMonto facturación\nCarta convenio · 19+', RED),
    ('Al cerrar', 'Diagnóstico *\nCausa raíz *\nTiempo resp. (auto)', AMBER),
]):
    x = 7.1 + i * 1.85
    rect(s, x, 4.45, 1.7, 1.5, BG_CARD, BORDER, True)
    rect(s, x, 4.45, 1.7, 0.04, c)
    txt(s, x + 0.1, 4.55, 1.5, 0.25, title, 10, True, CHARCOAL)
    multi(s, x + 0.1, 4.85, 1.5, 1, [(l, False, GRAY_DARK) for l in items.split('\n')], 8, GRAY_DARK, 1)

rect(s, 7.1, 6.15, 5.4, 0.4, RED_SOFT, radius=True)
txt(s, 7.3, 6.15, 5.2, 0.4, '⚠️  Diagnóstico y Causa raíz son OBLIGATORIOS para cerrar un ticket', 9, True, RED)

txt(s, 7.1, 6.7, 5.4, 0.3, '* Al marcar como Resuelto o Cerrado, el sistema los pide automáticamente', 8, False, GRAY)

# ═══════════════════════════════════════
# SLIDE 4: EQUIPO
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, 13.33, 7.5, OFFWHITE)
rect(s, 0, 0, 13.33, 0.06, BLUE)

txt(s, 0.8, 0.4, 1, 0.4, '03', 28, True, BLUE)
txt(s, 1.5, 0.45, 8, 0.4, '¿Qué debe hacer el equipo?', 24, True, CHARCOAL)
txt(s, 1.5, 0.85, 8, 0.3, 'Instrucciones para canalizar clientes al nuevo sistema', 11, False, GRAY_DARK)

# — Regla principal —
rect(s, 0.5, 1.35, 12.3, 1.1, GREEN_SOFT, GREEN, True)
txt(s, 0.9, 1.45, 6, 0.3, '✅  REGLA PRINCIPAL', 14, True, RGBColor(0x15,0x5E,0x2B))
txt(s, 0.9, 1.8, 7, 0.4, 'Cuando un cliente reporte una queja, oriéntalo a escribir "Hola" al WhatsApp de soporte.', 12, False, CHARCOAL)
rect(s, 9, 1.5, 3.5, 0.7, RGBColor(0x15,0x2E,0x1A), GREEN, True)
txt(s, 9, 1.5, 3.5, 0.7, '+52 449 156 1238', 20, True, GREEN, PP_ALIGN.CENTER)

# — Por área —
rect(s, 0.5, 2.7, 6, 3.5, WHITE, BORDER, True)
rect(s, 0.5, 2.7, 6, 0.06, BLUE)
txt(s, 0.8, 2.85, 4, 0.3, '🏢  Instrucciones por área', 13, True, CHARCOAL)

areas = [
    ('Administración', 'Compartir número al cliente que llama'),
    ('Cobranza', 'Compartir número en llamadas de seguimiento'),
    ('Operaciones', 'Compartir número al recibir queja en campo'),
    ('Ventas', 'Compartir guía visual a clientes cautivos'),
    ('Servicio al Cliente', 'Monitorear dashboard, atender chat, cerrar tickets'),
]
for i, (area, accion) in enumerate(areas):
    y = 3.35 + i * 0.5
    rect(s, 0.8, y, 0.06, 0.3, BLUE)
    txt(s, 1.0, y, 1.8, 0.3, area, 10, True, CHARCOAL)
    txt(s, 2.9, y, 3.3, 0.3, accion, 10, False, GRAY_DARK)

# — Puntos clave —
rect(s, 6.8, 2.7, 6, 2.3, WHITE, BORDER, True)
rect(s, 6.8, 2.7, 6, 0.06, AMBER)
txt(s, 7.1, 2.85, 5, 0.3, '💡  Puntos clave', 13, True, CHARCOAL)

puntos = [
    'Todos los reportes se canalizan por WhatsApp',
    'Folio automático con trazabilidad completa',
    'Agentes reciben notificación WhatsApp inmediata',
    'Tickets nunca se pierden (respaldados en Supabase)',
    'Al cerrar: diagnóstico + causa raíz obligatorios',
    'Encuesta CSAT (1-5) automática al cliente',
]
for i, p in enumerate(puntos):
    y = 3.3 + i * 0.27
    txt(s, 7.1, y, 5.5, 0.25, f'→  {p}', 9, False, CHARCOAL)

# — Links —
rect(s, 6.8, 5.2, 6, 2, WHITE, BORDER, True)
rect(s, 6.8, 5.2, 6, 0.06, PURPLE)
txt(s, 7.1, 5.35, 5, 0.3, '🔗  Links del sistema', 13, True, CHARCOAL)

links = [
    ('WhatsApp Bot', 'wa.me/524491561238', GREEN),
    ('Dashboard CX', 'cx-dashboard-54yw.onrender.com', BLUE),
    ('Chat Agentes', 'avantex-bot.onrender.com/dashboard', PURPLE),
    ('Guía pública', 'cx-dashboard-54yw.onrender.com/guia.html', AMBER),
]
for i, (name, url, c) in enumerate(links):
    y = 5.8 + i * 0.3
    rect(s, 7.1, y, 0.06, 0.2, c)
    txt(s, 7.3, y, 1.8, 0.2, name, 9, True, CHARCOAL)
    txt(s, 9.2, y, 3.3, 0.2, url, 9, False, BLUE)

# Footer
txt(s, 0.8, 6.85, 5, 0.25, 'Customer Experience  ·  Grupo Avantex  ·  Abril 2026', 9, False, GRAY)

# ═══════════════ GUARDAR ═══════════════
out = '/Users/diego/migracion-supabase/cx-dashboard/Avantex_Bot_WhatsApp_Dashboard_CX.pptx'
prs.save(out)
print(f'✅ PPT generado: {out}')
